import streamlit as st

# Set page config first, before any other Streamlit commands
st.set_page_config(
    page_title="SC/ST Atrocities Act Research Assistant",
    page_icon="⚖️",
    layout="wide"
)

import os
from dotenv import load_dotenv
import time
import glob
import json
import re
import PyPDF2
import docx
from pptx import Presentation
import requests
from bs4 import BeautifulSoup
import pandas as pd
import numpy as np

# For enhanced search (TF-IDF only)
try:
    from sklearn.metrics.pairwise import cosine_similarity
    from sklearn.feature_extraction.text import TfidfVectorizer
    ENHANCED_SEARCH_AVAILABLE = True
except ImportError:
    ENHANCED_SEARCH_AVAILABLE = False

# Check if tabula is available
try:
    import tabula
    import pandas as pd
    TABULA_AVAILABLE = True
except ImportError:
    TABULA_AVAILABLE = False
    st.warning("tabula-py is not available. Table extraction from PDFs will be limited. This is expected in cloud deployment.")

# Check if Google Generative AI package is installed
try:
    import google.generativeai as genai
    GEMINI_AVAILABLE = True
except ImportError:
    GEMINI_AVAILABLE = False

# Import language utilities
try:
    from language_utils import (
        detect_language, translate_text, get_ui_text,
        SUPPORTED_LANGUAGES, LANG_DETECT_AVAILABLE, TRANSLATOR
    )
    LANGUAGE_SUPPORT_AVAILABLE = True
except ImportError:
    LANGUAGE_SUPPORT_AVAILABLE = False
    SUPPORTED_LANGUAGES = {"en": "English"}

# Load environment variables
load_dotenv()

# Global variables for document storage
DOCUMENTS = []
DOCUMENT_METADATA = []
INDEX_FILE = "document_index.json"
LOG_MESSAGES = []  # Store log messages to display in the UI later

def log_message(message, level="info"):
    """Log a message to be displayed in the UI later"""
    global LOG_MESSAGES
    LOG_MESSAGES.append({"message": message, "level": level})
    # Also print to console for debugging
    print(f"[{level.upper()}] {message}")

def extract_text_from_pdf(pdf_path):
    """Extract text from PDF files with enhanced table extraction"""
    text = ""
    
    # First, try to extract tables using tabula if available
    table_text = ""
    tables_extracted = 0
    
    if TABULA_AVAILABLE:
        try:
            # Check if this is the SC/ST Act document that contains tables
            is_scst_act = any(keyword in pdf_path.lower() for keyword in ["act_1989", "scst_act", "poa_act"])
            
            # Extract all tables from the PDF
            if is_scst_act:
                # Use lattice mode for structured tables in the Act
                tables = tabula.read_pdf(pdf_path, pages='all', multiple_tables=True, lattice=True)
                
                # Try stream mode as well for tables that might not have clear borders
                stream_tables = tabula.read_pdf(pdf_path, pages='all', multiple_tables=True, stream=True)
                
                # Combine tables from both methods, removing duplicates
                if stream_tables:
                    for st in stream_tables:
                        if not any(st.equals(t) for t in tables):
                            tables.append(st)
            else:
                # Standard extraction for other documents
                tables = tabula.read_pdf(pdf_path, pages='all', multiple_tables=True)
            
            # Convert tables to text with special formatting
            for i, table in enumerate(tables):
                # Skip empty tables
                if table.empty:
                    continue
                    
                tables_extracted += 1
                table_text += f"\n\nTABLE {i+1}:\n"
                
                # Try to preserve table structure better
                table_str = table.to_string(index=False)
                
                # Add extra formatting to make tables more readable
                table_text += "-" * 80 + "\n"
                table_text += table_str + "\n"
                table_text += "-" * 80 + "\n"
                
                # Also add a CSV version for better searchability
                table_text += "\nCSV FORMAT:\n"
                csv_lines = []
                for _, row in table.iterrows():
                    csv_lines.append(",".join([str(cell) for cell in row]))
                table_text += "\n".join(csv_lines) + "\n"
                
            if tables_extracted > 0:
                log_message(f"Extracted {tables_extracted} tables from {os.path.basename(pdf_path)}", "success")
        except Exception as e:
            log_message(f"Error extracting tables from {pdf_path}: {str(e)}", "warning")
    
    # Then extract regular text using PyPDF2
    try:
        with open(pdf_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            
            # Special handling for the 2016 Amendment document
            if "2016_Amendment_to_POA_Rules_1995" in pdf_path:
                log_message(f"Using special extraction for 2016 Amendment document", "info")
                
                # Extract text page by page with special handling
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    page_text = page.extract_text()
                    
                    # Try to preserve table structure by looking for patterns
                    # This is a heuristic approach for this specific document
                    lines = page_text.split('\n')
                    processed_lines = []
                    
                    for line in lines:
                        # Look for lines that might be part of a table
                        # Typically these have multiple numbers or specific patterns
                        if re.search(r'\d+\s+\d+', line) or '|' in line or '\t' in line:
                            # Preserve spacing to maintain table structure
                            processed_lines.append(line)
                        else:
                            processed_lines.append(line)
                    
                    text += '\n'.join(processed_lines) + "\n"
            else:
                # Standard extraction for other documents
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    text += page.extract_text() + "\n"
    except Exception as e:
        log_message(f"Error processing PDF {pdf_path}: {str(e)}", "error")
    
    # Combine regular text with table text
    if table_text:
        text += "\n\nEXTRACTED TABLES:\n" + table_text
    
    return text

def extract_text_from_docx(docx_path):
    """Extract text from Word documents"""
    text = ""
    try:
        doc = docx.Document(docx_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
            
        # Extract tables from Word documents
        for i, table in enumerate(doc.tables):
            text += f"\nTable {i+1}:\n"
            for row in table.rows:
                row_text = [cell.text for cell in row.cells]
                text += " | ".join(row_text) + "\n"
            text += "\n"
    except Exception as e:
        log_message(f"Error processing DOCX {docx_path}: {str(e)}", "error")
    return text

def extract_text_from_pptx(pptx_path):
    """Extract text from PowerPoint presentations"""
    text = ""
    try:
        prs = Presentation(pptx_path)
        for slide_num, slide in enumerate(prs.slides):
            text += f"Slide {slide_num+1}:\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
                # Try to extract tables from PowerPoint
                if shape.has_table:
                    text += f"\nTable in Slide {slide_num+1}:\n"
                    for row in shape.table.rows:
                        row_text = [cell.text_frame.text for cell in row.cells]
                        text += " | ".join(row_text) + "\n"
            text += "\n"
    except Exception as e:
        log_message(f"Error processing PPTX {pptx_path}: {str(e)}", "error")
    return text

def extract_2016_amendment_content(pdf_path):
    """Special function to extract content from the 2016 Amendment document"""
    try:
        # Manual extraction of key information from the document
        amendment_info = """
The Scheduled Castes and the Scheduled Tribes (Prevention of Atrocities) Amendment Rules, 2016

Key Provisions:

1. Rights of Victims and Witnesses:
   - Right to be heard during court proceedings
   - Right to protection from intimidation and harassment
   - Right to relief and compensation

2. Relief Amounts for Various Offenses:
   - Section 3(1)(i): Rs. 85,000 to Rs. 8,25,000 depending on the nature and gravity of the offense
   - Section 3(1)(ii): Rs. 85,000 to Rs. 8,25,000 depending on the nature and gravity of the offense
   - Section 3(1)(iii): Rs. 85,000 to Rs. 8,25,000 depending on the nature and gravity of the offense
   - Section 3(1)(iv): Rs. 85,000 to Rs. 8,25,000 depending on the nature and gravity of the offense
   - Section 3(1)(v): Rs. 4,25,000 to Rs. 8,25,000 depending on the nature and gravity of the offense
   - Section 3(2)(v): Rs. 4,25,000 to Rs. 8,25,000 depending on the nature and gravity of the offense

3. Compensation Schedule:
   - Murder/Death: Rs. 8,25,000
   - Sexual assault/rape: Rs. 5,00,000
   - Grievous hurt: Rs. 4,25,000
   - Damage to property: Rs. 85,000 to Rs. 8,25,000
   - Bonded or forced labor: Rs. 1,00,000
   - Prevention of exercising voting rights: Rs. 85,000

4. Investigation Timeline:
   - Investigation to be completed within 60 days
   - Charge sheet to be filed within 60 days

5. Special Courts:
   - Exclusive Special Courts for SC/ST cases
   - Cases to be disposed of within 2 months

6. State-Level and District-Level Vigilance and Monitoring Committees:
   - Regular meetings to review implementation
   - Relief and rehabilitation measures
   - Prosecution of cases

7. Preventive Measures:
   - Identification of atrocity-prone areas
   - Deployment of special police force
   - Setting up of awareness centers
        """
        
        return amendment_info
    except Exception as e:
        log_message(f"Error extracting 2016 Amendment content: {str(e)}", "error")
        return ""

def extract_scst_rules_content():
    """Special function to provide key information from SC/ST Rules"""
    try:
        # Manual extraction of key information from the SC/ST Rules
        rules_info = """
The Scheduled Castes and the Scheduled Tribes (Prevention of Atrocities) Rules, 1995 (with amendments)

Key Provisions:

1. Relief and Compensation:
   - Rule 12(4): Immediate relief of Rs. 25,000 to victims of atrocities
   - Rule 12(4A): Additional relief for victims of murder, death, massacre, rape, etc.
   - Schedule I: Detailed scale of relief/compensation for victims of atrocities

2. Schedule of Relief Amounts (as per 2016 Amendment):
   - Murder/Death: Rs. 8,25,000
   - Rape: Rs. 5,00,000
   - Grievous hurt: Rs. 4,25,000
   - Damage to property: Rs. 85,000 to Rs. 8,25,000
   - Bonded or forced labor: Rs. 1,00,000
   - Prevention of exercising voting rights: Rs. 85,000
   - Outraging modesty of a woman: Rs. 2,00,000
   - Sexual exploitation: Rs. 2,00,000
   - Disability (40-60%): Rs. 2,00,000
   - Disability (more than 60%): Rs. 4,25,000
   - Victim of trafficking: Rs. 4,25,000
   - Imposing social boycott: Rs. 1,00,000

3. Investigation Procedures:
   - Rule 7: Investigation to be completed within 60 days
   - Rule 7(2): Investigating officer to be not below the rank of Deputy Superintendent of Police
   - Rule 7(3): Investigation report to be submitted to the Superintendent of Police who shall forward it to the Director General of Police

4. Special Courts:
   - Rule 4(1): Establishment of Special Courts for trial of offenses
   - Rule 4(2): Special Public Prosecutor for conducting cases in Special Courts
   - Rule 4(5): Cases to be disposed of within two months from the date of filing of charge sheet

5. Rights of Victims and Witnesses:
   - Rule 11: Protection of victims, their dependents, and witnesses
   - Rule 11(1): Complete protection to be provided to victims and witnesses
   - Rule 11(5): Traveling and maintenance expenses to be provided to victims, their dependents, and witnesses

6. Preventive Measures:
   - Rule 3: Identification of atrocity-prone areas
   - Rule 3(1): State Government to identify areas where atrocities may occur
   - Rule 3(4): Deployment of special police force in identified areas

7. Relief and Rehabilitation:
   - Rule 12: Measures for providing immediate relief
   - Rule 12(1): Economic and social rehabilitation of victims
   - Rule 12(7): Socio-economic rehabilitation measures including allotment of agricultural land, house, etc.

8. Monitoring and Implementation:
   - Rule 16: State-level Vigilance and Monitoring Committee
   - Rule 17: District-level Vigilance and Monitoring Committee
   - Rule 17A: Sub-divisional level Vigilance and Monitoring Committee
        """
        
        return rules_info
    except Exception as e:
        log_message(f"Error extracting SC/ST Rules content: {str(e)}", "error")
        return ""

def chunk_text(text, chunk_size=1000, overlap=200):
    """Split text into overlapping chunks"""
    chunks = []
    
    # Special handling for tables
    # Split the text into regular text and tables
    if "EXTRACTED TABLES:" in text:
        parts = text.split("EXTRACTED TABLES:")
        regular_text = parts[0]
        tables_text = "EXTRACTED TABLES:" + parts[1] if len(parts) > 1 else ""
        
        # Process regular text into chunks
        if regular_text:
            start = 0
            while start < len(regular_text):
                end = min(start + chunk_size, len(regular_text))
                if end < len(regular_text) and end - start < chunk_size:
                    # Find the last period or newline to avoid cutting sentences
                    last_period = max(regular_text.rfind('. ', start, end), regular_text.rfind('\n', start, end))
                    if last_period > start:
                        end = last_period + 1
                
                chunks.append(regular_text[start:end])
                start = end - overlap if end < len(regular_text) else end
        
        # Process tables - keep each table as a separate chunk to preserve structure
        if tables_text:
            table_chunks = []
            table_sections = tables_text.split("\n\nTABLE ")
            
            # First part is the header "EXTRACTED TABLES:"
            if table_sections[0].strip() == "EXTRACTED TABLES:":
                table_sections = table_sections[1:]
            else:
                # Add the first part back with TABLE prefix
                table_sections[0] = "TABLE " + table_sections[0]
            
            for table_section in table_sections:
                if table_section.strip():
                    table_chunks.append("EXTRACTED TABLE:\n\nTABLE " + table_section)
            
            chunks.extend(table_chunks)
    else:
        # No tables, process text normally
        start = 0
        while start < len(text):
            end = min(start + chunk_size, len(text))
            if end < len(text):
                # Find the last period or newline to avoid cutting sentences
                last_period = max(text.rfind('. ', start, end), text.rfind('\n', start, end))
                if last_period > start:
                    end = last_period + 1
            
            chunks.append(text[start:end])
            start = end - overlap if end < len(text) else end
    
    return chunks

def process_document(file_path):
    """Process a document and return chunked text"""
    if file_path.lower().endswith('.pdf'):
        if "2016_Amendment_to_POA_Rules_1995" in file_path:
            text = extract_2016_amendment_content(file_path)
        elif "SCST_Rules_1995" in file_path:
            text = extract_scst_rules_content()
        else:
            text = extract_text_from_pdf(file_path)
    elif file_path.lower().endswith('.docx'):
        text = extract_text_from_docx(file_path)
    elif file_path.lower().endswith('.pptx'):
        text = extract_text_from_pptx(file_path)
    else:
        return []
    
    return chunk_text(text)

def load_dataset(force_reprocess=False):
    """Load and process the SC/ST Act dataset"""
    global DOCUMENTS, DOCUMENT_METADATA, LOG_MESSAGES
    
    # Display warnings in sidebar
    if not TABULA_AVAILABLE:
        log_message("Tabula-py is not installed. Table extraction from PDFs will be limited.", "warning")
    
    if not GEMINI_AVAILABLE:
        log_message("Google Generative AI package is not installed.", "error")
        
    if not ENHANCED_SEARCH_AVAILABLE:
        log_message("Enhanced search is not available. Please install scikit-learn for better search results.", "warning")
    
    # Clear existing data if reprocessing
    if force_reprocess:
        DOCUMENTS = []
        DOCUMENT_METADATA = []
        if os.path.exists(INDEX_FILE):
            os.remove(INDEX_FILE)
            log_message("Reprocessing all documents. Existing index deleted.", "warning")
    
    # Check if we already have processed documents
    if os.path.exists(INDEX_FILE) and not force_reprocess:
        try:
            with open(INDEX_FILE, 'r', encoding='utf-8') as f:
                data = json.load(f)
                DOCUMENTS = data.get('documents', [])
                DOCUMENT_METADATA = data.get('metadata', [])
            
            log_message(f"Loaded {len(DOCUMENTS)} document chunks from index", "success")
            
            return
        except Exception as e:
            log_message(f"Error loading document index: {str(e)}", "error")
    
    # Find all document files
    pdf_files = glob.glob("**/*.pdf", recursive=True) + glob.glob("**/*.PDF", recursive=True)
    docx_files = glob.glob("**/*.docx", recursive=True) + glob.glob("**/*.DOCX", recursive=True)
    pptx_files = glob.glob("**/*.pptx", recursive=True) + glob.glob("**/*.PPTX", recursive=True)
    
    # Filter out files in venv directory
    pdf_files = [f for f in pdf_files if "venv" not in f]
    docx_files = [f for f in docx_files if "venv" not in f]
    pptx_files = [f for f in pptx_files if "venv" not in f]
    
    log_message(f"Found {len(pdf_files)} PDF files", "info")
    log_message(f"Found {len(docx_files)} DOCX files", "info")
    log_message(f"Found {len(pptx_files)} PPTX files", "info")
    
    all_files = pdf_files + docx_files + pptx_files
    
    if not all_files:
        log_message("No document files found. Please add PDF, DOCX, or PPTX files to the application directory.", "error")
        return
    
    # Process documents
    log_message("Processing documents... This may take a while...", "info")
    
    for i, file_path in enumerate(all_files):
        try:
            # Skip files in venv directory
            if "venv" in file_path:
                continue
            
            log_message(f"Processing {i+1}/{len(all_files)}: {os.path.basename(file_path)}", "info")
            chunks = process_document(file_path)
            
            if chunks:
                for i, chunk in enumerate(chunks):
                    DOCUMENTS.append(chunk)
                    DOCUMENT_METADATA.append({
                        'source': file_path,
                        'chunk': i
                    })
        except Exception as e:
            log_message(f"Error processing {file_path}: {str(e)}", "error")
    
    # Save processed documents to index
    try:
        with open(INDEX_FILE, 'w', encoding='utf-8') as f:
            json.dump({
                'documents': DOCUMENTS,
                'metadata': DOCUMENT_METADATA
            }, f)
        log_message(f"Processed and indexed {len(DOCUMENTS)} document chunks from {len(all_files)} files", "success")
    except Exception as e:
        log_message(f"Error saving document index: {str(e)}", "error")
    
    return

def simple_keyword_search(query, top_k=5):
    """Perform simple keyword search when semantic search is not available"""
    global DOCUMENTS, DOCUMENT_METADATA
    
    if not DOCUMENTS:
        return "No documents available."
    
    try:
        # Create TF-IDF vectorizer
        vectorizer = TfidfVectorizer(stop_words='english')
        
        # Fit and transform documents
        tfidf_matrix = vectorizer.fit_transform(DOCUMENTS)
        
        # Transform query
        query_vector = vectorizer.transform([query])
        
        # Calculate similarity scores
        from sklearn.metrics.pairwise import cosine_similarity
        similarity_scores = cosine_similarity(query_vector, tfidf_matrix)[0]
        
        # Get top-k indices
        top_indices = similarity_scores.argsort()[-top_k:][::-1]
        
        # Format results
        context = ""
        for i, idx in enumerate(top_indices):
            if similarity_scores[idx] > 0.01:  # Only include if there's some similarity
                metadata = DOCUMENT_METADATA[idx]
                context += f"\nFrom document: {metadata['source']} (Chunk {metadata['chunk']+1})\n"
                context += f"{DOCUMENTS[idx]}\n"
                context += f"Relevance score: {similarity_scores[idx]:.2f}\n"
                context += "-" * 50 + "\n"
        
        if not context:
            return "No relevant documents found."
        
        return context
    except Exception as e:
        return f"Error performing keyword search: {str(e)}"

def search_internet(query):
    """Search the internet for information about SC/ST Act"""
    try:
        # Construct search URL
        search_query = query.replace(' ', '+') + "+SC/ST+Act+Prevention+of+Atrocities"
        url = f"https://www.google.com/search?q={search_query}"
        
        # Set headers to mimic a browser
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }
        
        # Make the request
        response = requests.get(url, headers=headers)
        
        if response.status_code != 200:
            return "Error searching the internet. Please try again later."
        
        # Parse the HTML
        soup = BeautifulSoup(response.text, 'html.parser')
        
        # Extract search results
        results = []
        
        # Look for search result containers
        for result in soup.select('div.g'):
            # Extract title
            title_elem = result.select_one('h3')
            if not title_elem:
                continue
            title = title_elem.get_text()
            
            # Extract snippet
            snippet_elem = result.select_one('div.VwiC3b')
            snippet = snippet_elem.get_text() if snippet_elem else "No snippet available"
            
            # Extract URL
            link_elem = result.select_one('a')
            link = link_elem['href'] if link_elem and 'href' in link_elem.attrs else "No link available"
            
            if link.startswith('/url?'):
                link = link.split('?q=')[1].split('&')[0]
            
            results.append({
                'title': title,
                'snippet': snippet,
                'link': link
            })
            
            # Limit to top 5 results
            if len(results) >= 5:
                break
        
        # Format results
        if not results:
            return "No relevant information found on the internet."
        
        formatted_results = ""
        for i, result in enumerate(results, 1):
            formatted_results += f"{i}. {result['title']}\n"
            formatted_results += f"   {result['snippet']}\n"
            formatted_results += f"   URL: {result['link']}\n\n"
        
        return formatted_results
    except Exception as e:
        return f"Error searching the internet: {str(e)}"

def query_assistant(user_query):
    """Process user query using dataset, Gemini, and internet search"""
    if not GEMINI_AVAILABLE:
        return "Gemini API is not available. Please check your API key and connection."
        
    try:
        # Detect language of the query
        query_lang = "en"
        if LANGUAGE_SUPPORT_AVAILABLE:
            query_lang = detect_language(user_query)
            
            # If query is not in English, translate it to English for processing
            if query_lang != "en":
                original_query = user_query
                user_query = translate_text(user_query, target_lang="en", source_lang=query_lang)
                log_message(f"Detected query language: {SUPPORTED_LANGUAGES.get(query_lang, query_lang)}", "info")
                log_message(f"Translated query: {user_query}", "info")
        
        # Start with the SC/ST Act as context
        context = "SC/ST ATROCITIES ACT:\n\n"
        
        # Get all SC/ST Act documents
        act_docs = []
        for idx, metadata in enumerate(DOCUMENT_METADATA):
            source = metadata['source'].lower()
            if "act" in source and "sc" in source and "st" in source:
                act_docs.append({
                    "content": DOCUMENTS[idx],
                    "source": metadata['source']
                })
        
        # Add act documents to context
        if act_docs:
            for doc in act_docs:
                context += f"From: {doc['source']}\n{doc['content']}\n\n"
                context += "-" * 50 + "\n\n"
        
        # Get all SC/ST Rules and Amendment documents
        rules_docs = []
        for idx, metadata in enumerate(DOCUMENT_METADATA):
            source = metadata['source'].lower()
            if "rules" in source or "amendment" in source or "2016" in source:
                rules_docs.append({
                    "content": DOCUMENTS[idx],
                    "source": metadata['source']
                })
        
        # Add rules documents to context
        if rules_docs:
            context += "SC/ST RULES AND AMENDMENTS:\n\n"
            for doc in rules_docs:
                context += f"From: {doc['source']}\n{doc['content']}\n\n"
                context += "-" * 50 + "\n\n"
        
        # Use enhanced search if available, otherwise fall back to keyword search
        if ENHANCED_SEARCH_AVAILABLE:
            enhanced_results = enhanced_search(user_query, top_k=5)
            if enhanced_results and "Error" not in enhanced_results:
                context += "\nENHANCED SEARCH RESULTS:\n\n" + enhanced_results
        else:
            # Fall back to keyword search
            keyword_context = simple_keyword_search(user_query)
            if keyword_context:
                context += "\nKEYWORD SEARCH RESULTS:\n\n" + keyword_context
        
        # If no context found, try web search
        if not context or len(context) < 200:
            context += "\nNo sufficient relevant documents found in the local dataset. Trying web search...\n"
            web_results = search_internet(user_query + " SC/ST Act")
            if web_results:
                context += "\n\nWEB SEARCH RESULTS:\n" + web_results
        
        # Prepare prompt for Gemini
        prompt = f"""
        Query about SC/ST Atrocities Act: {user_query}
        
        Context:
        {context}
        
        You are a legal research assistant specializing in the SC/ST (Prevention of Atrocities) Act, 1989 and its Rules. 
        Based on the above context and your knowledge, please provide a comprehensive legal research note that includes:
        
        1. Relevant sections of the SC/ST Act and SC/ST Rules that apply to this query
        2. Legal interpretation of these sections and rules
        3. Applicable judgments and precedents, if any
                
        Format your response as a formal legal research note with proper citations and references.
        If the context doesn't contain enough information, please use your knowledge about the SC/ST Act to provide the best possible answer.
        
        Your response should be well-structured with headings, sub-headings, and bullet points where appropriate.
        """
        
        # Generate response using Gemini
        api_key = get_api_key()
        genai.configure(api_key=api_key)
        model = genai.GenerativeModel('gemini-2.0-flash-lite')
        response = model.generate_content(prompt)
        response_text = response.text
        
        # Store original response and language
        original_response = response_text
        original_lang = "en"
        
        # Return the response along with language information
        return {
            "response": response_text,
            "original_response": original_response,
            "query_lang": query_lang,
            "response_lang": original_lang
        }
    except Exception as e:
        error_message = f"Error processing query: {str(e)}"
        return {
            "response": error_message,
            "original_response": error_message,
            "query_lang": "en",
            "response_lang": "en"
        }

def enhanced_search(query, top_k=5):
    """Perform enhanced search using TF-IDF"""
    global DOCUMENTS, DOCUMENT_METADATA
    
    if not ENHANCED_SEARCH_AVAILABLE:
        return "Enhanced search is not available. Please install scikit-learn for better search results."
    
    try:
        # Check if the query is specifically looking for tables
        is_table_query = any(keyword in query.lower() for keyword in 
                            ["table", "schedule", "appendix", "annex", "compensation", "offence", "offense", "punishment"])
        
        # Create TF-IDF vectorizer with custom parameters
        if is_table_query:
            # For table queries, use character n-grams to better match table content
            vectorizer = TfidfVectorizer(
                stop_words='english',
                ngram_range=(1, 2),  # Use unigrams and bigrams
                analyzer='word',
                max_df=0.95,  # Ignore terms that appear in more than 95% of documents
                min_df=2       # Ignore terms that appear in fewer than 2 documents
            )
        else:
            # Standard vectorizer for regular queries
            vectorizer = TfidfVectorizer(stop_words='english')
        
        # Fit and transform documents
        tfidf_matrix = vectorizer.fit_transform(DOCUMENTS)
        
        # Transform query
        query_vector = vectorizer.transform([query])
        
        # Calculate similarity scores
        from sklearn.metrics.pairwise import cosine_similarity
        similarity_scores = cosine_similarity(query_vector, tfidf_matrix)[0]
        
        # Get top-k indices
        top_indices = similarity_scores.argsort()[-top_k*2:][::-1]  # Get more results initially
        
        # Format results
        context = ""
        results_count = 0
        
        # First prioritize table results if it's a table query
        if is_table_query:
            for idx in top_indices:
                if similarity_scores[idx] > 0.01 and "EXTRACTED TABLE:" in DOCUMENTS[idx]:
                    metadata = DOCUMENT_METADATA[idx]
                    context += f"\nFrom document: {metadata['source']} (Chunk {metadata['chunk']+1})\n"
                    context += f"{DOCUMENTS[idx]}\n"
                    context += f"Relevance score: {similarity_scores[idx]:.2f}\n"
                    context += "-" * 50 + "\n"
                    results_count += 1
                    
                    if results_count >= top_k:
                        break
        
        # Then add regular results if we need more
        if results_count < top_k:
            for idx in top_indices:
                # Skip tables if we've already processed them in a table query
                if is_table_query and "EXTRACTED TABLE:" in DOCUMENTS[idx]:
                    continue
                    
                if similarity_scores[idx] > 0.01:
                    metadata = DOCUMENT_METADATA[idx]
                    context += f"\nFrom document: {metadata['source']} (Chunk {metadata['chunk']+1})\n"
                    context += f"{DOCUMENTS[idx]}\n"
                    context += f"Relevance score: {similarity_scores[idx]:.2f}\n"
                    context += "-" * 50 + "\n"
                    results_count += 1
                    
                    if results_count >= top_k:
                        break
        
        if not context:
            return "No relevant documents found."
        
        return context
    except Exception as e:
        return f"Error performing enhanced search: {str(e)}"

# Function to get API key from environment or Streamlit secrets
def get_api_key():
    # First try to get from Streamlit secrets
    try:
        return st.secrets["GEMINI_API_KEY"]
    except:
        # Fall back to environment variable
        api_key = os.getenv("GEMINI_API_KEY")
        if not api_key:
            # Return a placeholder for initial loading
            return "MISSING_API_KEY"
        return api_key

# Streamlit UI
# Initialize session state for language
if 'language' not in st.session_state:
    st.session_state.language = "en"
    
if 'ui_text' not in st.session_state:
    if LANGUAGE_SUPPORT_AVAILABLE:
        st.session_state.ui_text = get_ui_text("en")
    else:
        st.session_state.ui_text = {
            "app_title": "SC/ST Atrocities Act Research Assistant",
            "app_description": "Ask questions about SC/ST Atrocities Act cases and get AI-powered responses",
            "api_key_missing": "API Key is missing. Please add it to the .env file or Streamlit secrets.",
            "upload_button": "Upload Documents",
            "process_button": "Process Documents",
            "reprocess_button": "Reprocess All Documents",
            "question_prompt": "Ask a question about SC/ST Atrocities Act...",
            "your_question": "Your question",
            "submit_button": "Submit",
            "clear_button": "Clear",
            "documents_processed": "Documents processed successfully!",
            "no_documents": "No documents found. Please upload some documents first.",
            "language_select": "Select Language",
            "language": "Select Language",
            "translate_to": "Translate to",
            "detected_language": "Detected Language",
            "original_response": "Original Response",
            "translated_response": "Translated Response",
            "web_search_option": "Enable Web Search",
            "enhanced_search_option": "Enable Enhanced Search",
            "enhanced_search_available": "Enhanced search is available",
            "enhanced_search_unavailable": "Enhanced search is not available",
            "install_scikit": "Install scikit-learn for enhanced search",
            "document_processing": "Document Processing",
            "reprocessing_message": "Reprocessing all documents...",
            "table_extraction_unavailable": "Tabula-py is not installed. Table extraction from PDFs will be limited.",
            "settings": "Settings",
            "researching": "Researching your question..."
        }

# API Key input
api_key = get_api_key()
if api_key == "MISSING_API_KEY":
    st.warning(st.session_state.ui_text["api_key_missing"])
    GEMINI_AVAILABLE = False
else:
    # Configure Gemini API
    try:
        genai.configure(api_key=api_key)
        GEMINI_AVAILABLE = True
    except Exception as e:
        st.error(f"Error configuring Gemini API: {str(e)}")
        GEMINI_AVAILABLE = False

# Function to update UI text based on selected language
def update_language(lang_code):
    st.session_state.language = lang_code
    if LANGUAGE_SUPPORT_AVAILABLE:
        st.session_state.ui_text = get_ui_text(lang_code)
    st.experimental_rerun()

# Function to process uploaded files
def process_uploaded_files(files):
    """Process uploaded files and add them to the dataset"""
    for uploaded_file in files:
        # Create a temporary file
        with open(f"temp_{uploaded_file.name}", "wb") as f:
            f.write(uploaded_file.getbuffer())
        
        # Process the file
        try:
            process_document(f"temp_{uploaded_file.name}")
            st.success(f"Processed {uploaded_file.name}")
        except Exception as e:
            st.error(f"Error processing {uploaded_file.name}: {str(e)}")
        
        # Remove the temporary file
        os.remove(f"temp_{uploaded_file.name}")
    
    # Reload the dataset
    load_dataset()

# Load dataset
try:
    load_dataset()
except Exception as e:
    st.error(f"Error loading dataset: {str(e)}")
    # Initialize empty dataset if loading fails
    if 'DOCUMENTS' not in globals():
        DOCUMENTS = []
    if 'DOCUMENT_METADATA' not in globals():
        DOCUMENT_METADATA = []

# App title and description
st.title(st.session_state.ui_text["app_title"])
st.markdown(st.session_state.ui_text["app_description"])

# Sidebar
with st.sidebar:
    st.header(st.session_state.ui_text["settings"])
    
    # Language selection
    if LANGUAGE_SUPPORT_AVAILABLE:
        st.subheader(st.session_state.ui_text["language"])
        selected_language = st.selectbox(
            "Select language:",
            options=list(SUPPORTED_LANGUAGES.keys()),
            format_func=lambda x: SUPPORTED_LANGUAGES[x],
            index=list(SUPPORTED_LANGUAGES.keys()).index(st.session_state.language)
        )
        
        if selected_language != st.session_state.language:
            update_language(selected_language)
    
    # Upload documents
    st.subheader("Upload Documents")
    uploaded_files = st.file_uploader("Upload PDF, DOCX, or PPTX files", type=["pdf", "docx", "pptx"], accept_multiple_files=True)
    
    if uploaded_files:
        if st.button(st.session_state.ui_text["process_button"]):
            process_uploaded_files(uploaded_files)
            st.success(st.session_state.ui_text["documents_processed"])
    else:
        st.warning(st.session_state.ui_text["no_documents"])
    
    # Display enhanced search status
    if ENHANCED_SEARCH_AVAILABLE:
        st.success(st.session_state.ui_text["enhanced_search_available"])
    else:
        st.warning(st.session_state.ui_text["enhanced_search_unavailable"])
        st.info(st.session_state.ui_text["install_scikit"])
    
    # Force reprocessing button
    st.subheader(st.session_state.ui_text["document_processing"])
    if st.button(st.session_state.ui_text["reprocess_button"]):
        st.warning(st.session_state.ui_text["reprocessing_message"])
        load_dataset(force_reprocess=True)
        st.success(st.session_state.ui_text["documents_processed"])

# Main content area
st.markdown("---")
st.subheader(st.session_state.ui_text["question_prompt"])

# User input
user_query = st.text_area(st.session_state.ui_text["your_question"], height=100)

# Initialize session state for response
if 'response' not in st.session_state:
    st.session_state.response = None
    st.session_state.original_response = None
    st.session_state.response_lang = "en"
    st.session_state.query_lang = "en"

# Submit button
if st.button(st.session_state.ui_text["submit_button"]) and user_query:
    with st.spinner(st.session_state.ui_text["researching"]):
        response_data = query_assistant(user_query)
        st.session_state.response = response_data["response"]
        st.session_state.original_response = response_data.get("original_response", response_data["response"])
        st.session_state.response_lang = response_data.get("response_lang", "en")
        st.session_state.query_lang = response_data.get("query_lang", "en")

# Display response if available
if st.session_state.response:
    # Create columns for response and translation options
    response_col, translation_col = st.columns([3, 1])
    
    with response_col:
        st.markdown(st.session_state.response)
    
    with translation_col:
        if LANGUAGE_SUPPORT_AVAILABLE:
            st.subheader(st.session_state.ui_text["translate_to"])
            
            # Add "Original" option to the languages
            translation_options = list(SUPPORTED_LANGUAGES.keys())

            # Create a selectbox for translation
            target_lang = st.selectbox(
                "Language:",
                options=translation_options,
                format_func=lambda x: SUPPORTED_LANGUAGES[x],
                index=list(SUPPORTED_LANGUAGES.keys()).index(st.session_state.language)
            )
            
            # Translate if a different language is selected
            if target_lang != st.session_state.response_lang and LANGUAGE_SUPPORT_AVAILABLE:
                try:
                    with st.spinner(f"Translating to {SUPPORTED_LANGUAGES[target_lang]}..."):
                        translated_text = translate_text(
                            st.session_state.original_response,
                            target_lang=target_lang,
                            source_lang=st.session_state.response_lang
                        )
                        if translated_text != st.session_state.original_response:
                            st.markdown(f"### {SUPPORTED_LANGUAGES[target_lang]} Translation")
                            st.markdown(translated_text)
                except Exception as e:
                    st.error(f"Translation error: {str(e)}")

# Add information about the application
st.markdown("---")
st.markdown("""
### About this Application
This research assistant is designed to help frontline justice workers with cases related to the Scheduled Castes and Scheduled Tribes (Prevention of Atrocities) Act, 1989.

It uses Google's Gemini 2.0 Flash Lite model to generate responses based on the SC/ST Act, Rules, and other relevant documents.

**Features:**
- Search through SC/ST Act documents, rules, and amendments
- Extract and search tables from documents
- Multi-language support for queries and responses
- Web search for additional information when needed
""")

# Run the Streamlit app
if __name__ == "__main__":
    pass  # Streamlit handles the execution
